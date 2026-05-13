import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
import os
import requests
import re
import io
import cv2
import copy
import time
from datetime import datetime
import streamlit as st

# --- STREAMLIT SETTINGS (Neon Green & Black Theme) ---
st.set_page_config(
    page_title="Auto Generate Report Observation", 
    layout="centered",
    initial_sidebar_state="collapsed"
)

# Custom CSS
st.markdown("""
    <style>
    .stApp { background-color: #0B0E0F; }
    .custom-header {
        background-color: #111516;
        border-bottom: 2px solid #00FF66;
        padding: 15px;
        border-radius: 8px 8px 0px 0px;
        display: flex;
        justify-content: space-between;
        align-items: center;
        margin-bottom: 20px;
    }
    .header-title {
        color: #00FF66 !important;
        font-family: 'Courier New', Courier, monospace;
        font-size: 22px;
        font-weight: bold;
        letter-spacing: 2px;
        margin: 0;
    }
    label, .stMarkdown p {
        color: #888888 !important;
        font-family: 'Courier New', Courier, monospace !important;
        text-transform: uppercase;
        font-size: 12px !important;
        font-weight: bold;
    }
    .stTextInput>div>div>input, .stSelectbox>div>div>div, .stDateInput>div>div>input {
        background-color: #111516 !important;
        color: #00FF66 !important;
        border: 1px solid #333333 !important;
        font-family: 'Courier New', Courier, monospace !important;
        border-radius: 4px !important;
    }
    .stAlert {
        background-color: #111516 !important;
        color: #00FF66 !important;
        border: 1px solid #00FF66 !important;
    }
    .stButton>button {
        background-color: transparent !important;
        color: #00FF66 !important;
        border: 2px solid #00FF66 !important;
        font-family: 'Courier New', Courier, monospace !important;
        font-weight: bold !important;
        transition: 0.3s;
    }
    .stButton>button:hover {
        background-color: #00FF66 !important;
        color: #000000 !important;
    }
    hr { border-color: #1A1E1F !important; }
    .stProgress > div > div > div > div { background-color: #00FF66 !important; }
    </style>
    """, unsafe_allow_html=True)

# --- HEADER ---
st.markdown("""
    <div class="custom-header">
        <p class="header-title">OBSERVATION DETAILS</p>
        <span style="color: #FF3366; font-weight: bold; font-family: sans-serif; cursor: pointer;">✕</span>
    </div>
    """, unsafe_allow_html=True)

# --- CORE LOGIC ---
def download_and_insert_media(slide, link_data, left_inch, top_inch, width_inch, is_video_slide=False):
    if not isinstance(link_data, str) or "drive.google.com" not in link_data:
        return
    video_folder = "downloaded_videos"
    if not os.path.exists(video_folder): os.makedirs(video_folder)
    links = [l.strip() for l in link_data.split(';') if l.strip()]

    for index, clean_link in enumerate(links):
        file_id_match = re.search(r'[-\w]{25,}', clean_link)
        if file_id_match:
            file_id = file_id_match.group(0)
            session = requests.Session()
            download_url = "https://docs.google.com/uc?export=download"
            params = {'id': file_id, 'confirm': 't'}
            current_left = 2.0 if is_video_slide else left_inch
            current_top = 1.5 if is_video_slide else top_inch + (index * 2.1)
            current_width = 6.0 if is_video_slide else width_inch

            try:
                response = session.get(download_url, params=params, stream=True, timeout=30)
                if response.status_code == 200:
                    content_type = response.headers.get('Content-Type', '')
                    if 'image' in content_type and not is_video_slide:
                        image_data = io.BytesIO(response.content)
                        slide.shapes.add_picture(image_data, Inches(current_left), Inches(current_top), width=Inches(current_width))
                    elif ('video' in content_type or 'octet-stream' in content_type) and is_video_slide:
                        video_path = os.path.join(video_folder, f"video_{file_id}.mp4")
                        thumb_path = os.path.join(video_folder, f"thumb_{file_id}.jpg")
                        with open(video_path, 'wb') as f:
                            for chunk in response.iter_content(chunk_size=8192): f.write(chunk)
                        vidcap = cv2.VideoCapture(video_path)
                        vidcap.set(cv2.CAP_PROP_POS_MSEC, 12000)
                        success, image = vidcap.read()
                        if not success:
                            vidcap.set(cv2.CAP_PROP_POS_MSEC, 0)
                            success, image = vidcap.read()
                        if success: cv2.imwrite(thumb_path, image)
                        vidcap.release()
                        slide.shapes.add_movie(video_path, Inches(current_left), Inches(current_top), width=Inches(current_width), height=Inches(current_width * 0.56), poster_frame_image=thumb_path if os.path.exists(thumb_path) else None, mime_type='video/mp4')
            except Exception: pass

def create_custom_slide(pres, slide_template):
    blank_layout = pres.slide_layouts[6]
    new_slide = pres.slides.add_slide(blank_layout)
    for shp in list(new_slide.shapes): new_slide.shapes._spTree.remove(shp.element)
    footer_threshold = pres.slide_height * 0.85
    for shape in slide_template.shapes:
        if shape.is_placeholder or shape.top > footer_threshold: continue
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            img_stream = io.BytesIO(shape.image.blob)
            new_slide.shapes.add_picture(img_stream, shape.left, shape.top, shape.width, shape.height)
        else:
            new_el = copy.deepcopy(shape.element)
            new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
    return new_slide

# --- UI INPUTS ---
TEMPLATE_FILENAME = "template.pptx"
template_exists = os.path.exists(TEMPLATE_FILENAME)
if not template_exists: st.error(f"STATUS: '{TEMPLATE_FILENAME}' NOT FOUND")
else: st.success(f"STATUS: TEMPLATE LOADED")

col1, col2 = st.columns(2)
with col1:
    report_title = st.text_input("REPORT TITLE", placeholder="Input Title...")
    obs_month = st.selectbox("OBSERVATION MONTH", ["January", "February", "March", "April", "May", "June", "July", "August", "September", "October", "November", "December"])
    start_date = st.date_input("START DATE", value=datetime(2026, 4, 27))
with col2:
    selected_depot = st.selectbox("DEPOT LOCATION", ["MRT Jinjang", "Shah Alam", "Cheras Selatan", "Batu Caves", "MRT Kajang", "MRT Sungai Buloh", "MRT Serdang"])
    selected_siri = st.selectbox("SIRI NUMBER", [f"OE/SQI/CR/VO/{str(i).zfill(3)}/2026" for i in range(1, 101)])
    end_date = st.date_input("END DATE", value=datetime(2026, 4, 30))

st.divider()

# GOOGLE SHEET SOURCE
SHEET_ID = "1qlPsPPRKMTfoyMN0MmzK3Hu9wxiBFjYIX6IFMbriZmo"
GID = "40024720"
CSV_URL = f"https://docs.google.com/spreadsheets/d/{SHEET_ID}/export?format=csv&gid={GID}"

btn_col1, btn_col2 = st.columns([3, 1])
with btn_col1: generate_btn = st.button("RUN GENERATOR", use_container_width=True)
with btn_col2:
    if st.button("REFRESH", use_container_width=True):
        st.cache_data.clear()
        st.rerun()

# --- EXECUTION ---
if generate_btn:
    if not template_exists:
        st.error("ERROR: TEMPLATE MISSING")
    else:
        try:
            start_time = time.time()
            df = pd.read_csv(CSV_URL)
            
            # Convert first column to datetime objects
            df.iloc[:, 0] = pd.to_datetime(df.iloc[:, 0], errors='coerce')
            
            # Filtering
            mask = (df.iloc[:, 0].dt.date >= start_date) & \
                   (df.iloc[:, 0].dt.date <= end_date) & \
                   (df.iloc[:, 3].astype(str).str.strip() == selected_depot)

            filtered_data = df.loc[mask].copy()

            if filtered_data.empty:
                st.warning(f"NO RECORDS FOUND FOR {selected_depot}")
            else:
                prs = Presentation(TEMPLATE_FILENAME)
                slide1_template, slide2_template, slide3_template = prs.slides[0], prs.slides[1], prs.slides[2]
                slide6_template = prs.slides[5] if len(prs.slides) >= 6 else None

                processed_count, total = 0, len(filtered_data)
                progress_bar, status_text = st.progress(0), st.empty()
                summary_list = []

                for i in range(len(filtered_data)):
                    # Extract single row
                    row = filtered_data.iloc[i]
                    
                    if str(row.iloc[8]).strip().lower() == "yes": continue
                    summary_list.append(row)
                    
                    new_slide = create_custom_slide(prs, slide1_template)
                    
                    # Extract single timestamp value correctly
                    dt_val = row.iloc[0]
                    date_str = dt_val.strftime('%d/%m/%Y') if pd.notnull(dt_val) else "N/A"
                    time_str = dt_val.strftime('%H:%M:%S') if pd.notnull(dt_val) else "N/A"

                    c_i, c_j = str(row.iloc[8]).lower(), str(row.iloc[9]).lower()
                    pemerhatian = ("1. Pelanggaran Had Laju Hentian: BC memintas hentian dengan kelajuan melebihi 25 km/j.\n" if c_i == "no" else "") + \
                                  ("2. Kapten Bas tidak memandu / menggunakan lorong kiri" if c_j == "no" else "")
                    
                    cadangan = ""
                    if c_i == "no" and c_j == "no": cadangan = "1. Memberi peringatan/kaunseling kepada Kapten Bas memperlahankan bas and keperluan berada di lorong kiri."
                    elif c_i == "no": cadangan = "1. Memberi peringatan/kaunseling kepada Kapten Bas memperlahankan bas di setiap hentian bas."
                    elif c_j == "no": cadangan = "2. Memberi peringatan kepada Kapten Bas mengenai keperluan berada di lorong kiri."

                    # Create strings for replacement
                    reps = {
                        "Tarikh pemerhatian :": f"Tarikh pemerhatian : {date_str}", "Nombor Bas :": f"Nombor Bas : {str(row.iloc[6])}",
                        "Laluan pemerhatian :": f"Laluan pemerhatian : {str(row.iloc[4])}", "Masa :": f"Masa : {time_str}",
                        "Lokasi / Hentian :": f"Lokasi / Hentian : {str(row.iloc[5])}", "Nama Kapten Bas :": f"Nama Kapten Bas : {str(row.iloc[32])}",
                        "ID Kapten Bas :": f"ID Kapten Bas : {str(row.iloc[31])}", "Kelajuan Dipandu :": f"Kelajuan Dipandu : {str(row.iloc[30])} Km/h",
                        "Nama PIC :": f"Nama PIC : {str(row.iloc[2])}", "Pemerhatian Pemanduan Kapten Bas :": f"Pemerhatian Pemanduan Kapten Bas :\n{pemerhatian}",
                        "Cadangan:": f"Cadangan:\n{cadangan}"
                    }

                    for shape in new_slide.shapes:
                        if shape.has_table:
                            for r in shape.table.rows:
                                for cell in r.cells:
                                    for para in cell.text_frame.paragraphs:
                                        for k, v in reps.items():
                                            if k in para.text:
                                                para.text = para.text.replace(k, v)
                                                for run in para.runs: run.font.size = Pt(10)

                    download_and_insert_media(new_slide, row.iloc[26], 0.6, 2.1, 3.8)
                    v_slide = create_custom_slide(prs, slide2_template)
                    download_and_insert_media(v_slide, row.iloc[26], 0, 0, 0, True)

                    processed_count += 1
                    progress_bar.progress(processed_count / total)
                    status_text.text(f"COMPILING... {processed_count}/{total}")

                progress_bar.progress(1.0)
                status_text.text(f"DONE!")

                if summary_list:
                    new_summary_slide = create_custom_slide(prs, slide3_template)
                    orig = next((s for s in new_summary_slide.shapes if s.has_table), None)
                    if orig:
                        new_table = new_summary_slide.shapes.add_table(len(summary_list)+1, 6, Inches(0.5), Inches(1.5), Inches(9.0), orig.height).table
                        new_summary_slide.shapes._spTree.remove(orig.element)
                        headers = ["Depoh", "Laluan", "Nombor Bas", "Hentian Bas", "Pengesahan", "Status"]
                        for i, h in enumerate(headers): 
                            new_table.rows[0].cells[i].text = h
                            for p in new_table.rows[0].cells[i].text_frame.paragraphs:
                                for r in p.runs: r.font.size, r.font.bold = Pt(10), True
                        for idx, s_row in enumerate(summary_list):
                            tr = new_table.rows[idx+1]
                            tr.cells[0].text, tr.cells[1].text, tr.cells[2].text, tr.cells[3].text = str(s_row.iloc[3]), str(s_row.iloc[4]), str(s_row.iloc[6]), str(s_row.iloc[5])
                            tr.cells[4].text = f"ID: {s_row.iloc[31]}\nNama: {s_row.iloc[32]}\nLaju: {s_row.iloc[30]} Km/h\nMasa: {pd.to_datetime(s_row.iloc[0]).strftime('%d/%m/%Y %H:%M:%S')}"
                            tr.cells[5].text = "Tidak Mematuhi"
                            for cell in tr.cells:
                                for p in cell.text_frame.paragraphs:
                                    for r in p.runs: r.font.size, r.font.name = Pt(8), "Arial"

                if slide6_template: create_custom_slide(prs, slide6_template)

                # Reordering
                xml_slides = prs.slides._sldIdLst
                for _ in range(3): xml_slides.remove(xml_slides[0])
                if len(prs.slides) >= 3:
                    summ_el = xml_slides[len(xml_slides)-2]
                    xml_slides.remove(summ_el); xml_slides.insert(2, summ_el)
                    s2_el = xml_slides[1]
                    xml_slides.remove(s2_el); xml_slides.insert(0, s2_el)
                if len(xml_slides) >= 4: xml_slides.remove(xml_slides[3])

                full_title = f"Central Region {obs_month} {datetime.now().year}"
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            if "Januari-February 2026" in shape.text_frame.text:
                                for p in shape.text_frame.paragraphs:
                                    if "Januari-February 2026" in p.text:
                                        p.text = p.text.replace("Januari-February 2026", full_title)
                                        for r in p.runs: r.font.name, r.font.size, r.font.color.rgb = 'Arial', Pt(20), RGBColor(0,32,96)
                            if "OE/SQI/CR/VO/001/2026" in shape.text_frame.text:
                                for p in shape.text_frame.paragraphs:
                                    if "OE/SQI/CR/VO/001/2026" in p.text:
                                        p.text = p.text.replace("OE/SQI/CR/VO/001/2026", selected_siri)
                                        for r in p.runs: r.font.size = Pt(12)

                ppt_out = io.BytesIO()
                prs.save(ppt_out)
                ppt_out.seek(0)
                st.success(f"COMPLETE: {int((time.time()-start_time)//60)}M {int((time.time()-start_time)%60)}S")
                st.download_button("DOWNLOAD GENERATED PPTX", ppt_out, f"{report_title or 'Report'}.pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation", use_container_width=True)
        except Exception as e: st.error(f"SYSTEM ERROR: {str(e)}")
